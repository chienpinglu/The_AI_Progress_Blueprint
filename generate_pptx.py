from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Slide definitions
slides_data = [
    {
        "title": "The Architecture of AI Progress",
        "bullets": [
            "Synthesizing the Formal Foundations and Empirical Drivers of Artificial Intelligence",
            "A Unified Theoretical Framework: Presenting a cohesive research program that bridges theoretical limits and engineering realities",
            "Two Connected Pairs:",
            "  Pair 1 (The Formal Foundation): Intelligence defined by its outer boundary of computability and inner boundary of polynomial-time feasibility",
            "  Pair 2 (The Practical Bridge): The empirical driver of scaling laws and the architectural consequences of hardware allocation"
        ]
    },
    {
        "title": "The Discourse Gap: Why This Program Exists",
        "bullets": [
            "The Hype: Empirical AI discourse often focuses on capability scaling and scaling laws without a rigorous formal theory of progress",
            "The Limits: Theoretical discourse frequently emphasizes strict mathematical limits without explaining how relentless practical progress continues to appear",
            "The Missing Link: We require a single, unified model that explains how formal boundaries, empirical drivers, and system-level architectural consequences interact"
        ]
    },
    {
        "title": "The Two-Pair Structure of AI Progress",
        "bullets": [
            "Pair 1: The Formal Foundation",
            "  Paper 1 (Computability): Intelligence as an open-ended limiting process at the outer boundary of what can be mechanized",
            "  Paper 2 (Complexity): Intelligence as a staged bounded ascent at the inner boundary of polynomial-time reachability",
            "Pair 2: The Empirical and Architectural Bridge",
            "  Paper 3 (Scaling Laws): How abstract scaling drives progress and enforces a race for hidden efficiency",
            "  Paper 4 (Modernized Amdahl): How the pressure of scaling reshapes the hardware allocation between specialized and programmable compute"
        ]
    },
    {
        "title": "Paper 1: The Computability Boundary",
        "bullets": [
            "Beyond a Fixed Closure: Intelligence is not one finitely specifiable, fixed rule-governed procedure",
            "The Limiting-Process Structure: When a process is mechanized, the higher-order task of guiding, specifying, and evaluating that mechanization is displaced upward to a new level",
            "Local Mechanization vs. Open Boundary: Every established, mechanized computational layer possesses a canonical, relatively uncomputable open boundary"
        ]
    },
    {
        "title": "Paper 1: The Mechanics of Ascent",
        "bullets": [
            "[Visual Concept: A blueprint-style diagram showing a foundational layer A_n, ascending through an intermediate step, toward an open ceiling A'_n]",
            "The Turing Jump as a Ceiling: The absolute boundary of any mechanized layer is its Turing jump (A')",
            "Progressive Layers: Real progress occurs through intermediate strengthenings (A <_T B <_T A'), mechanizing steps strictly below the absolute jump",
            "Non-Capture: No single settled layer captures the entire open-ended process of intelligence"
        ]
    },
    {
        "title": "Paper 2: The Feasibility Boundary",
        "bullets": [
            "Redefining Feasibility: The practical reach of intelligence is not a binary “P = NP or nothing” scenario",
            "Expanding Reachability: Progress manifests as progressively wider polynomial-time reach across difficult search tasks",
            "Staged Bounded Ascent: Moving upward by internalizing richer structures of guided search while remaining feasible under strict polynomial-time constraints"
        ]
    },
    {
        "title": "Paper 2: The Polynomial Hierarchy as Limiting Process",
        "bullets": [
            "[Visual Concept: A finite-depth ladder expanding upward, showing P→NP→PH]",
            "Bounded Verification Games: The Polynomial Hierarchy provides a finite-depth ladder of structured proposal, challenge, and repair",
            "Internalizing Search: Moving up the hierarchy allows a system to absorb what was previously external search into a stronger—but still explicitly bounded—verification framework"
        ]
    },
    {
        "title": "Synthesis of Pair 1: The Map of Formal Boundaries",
        "bullets": [
            "[Visual Concept: A nested boundary map showing Well-defined tasks → Computable tasks → Polynomial-time tasks]",
            "The Outer Boundary (Computability): The divide between computable tasks and uncomputable, open-ended well-defined tasks",
            "The Inner Boundary (Feasibility): The divide between what is merely computable in principle and what is polynomial-time reachable in practice"
        ]
    },
    {
        "title": "The Empirical Gap: Why Pair 1 is Not Enough",
        "bullets": [
            "The Transition: Pair 1 provides the formal structure, explaining exactly what progress means mathematically",
            "The Missing Engine: Formal limits alone do not explain why progress continually appears in real-world systems",
            "We must bridge the gap between formal structure and empirical hardware realities"
        ]
    },
    {
        "title": "Paper 3: Scaling Laws and Logical Compute",
        "bullets": [
            "The Power of Abstraction: Scaling laws are unreasonably effective precisely because they abstract over realization details",
            "Logical vs. Physical Compute: The \"compute\" variable actually represents logical compute (abstract model-side work), independent of the efficiency of the physical hardware",
            "The Real Meaning of Diminishing Returns: A flattening scaling curve does not just mean slower progress; it represents a rapidly escalating physical and operational burden to achieve the next step of value"
        ]
    },
    {
        "title": "Paper 3: The Dynamic Implication of Hidden Efficiency",
        "bullets": [
            "[Visual Concept: A pressure-conversion loop showing logical compute demand driving efficiency optimization across the stack]",
            "The Efficiency Race: Progress persists because of repeated, compounding efficiency gains in algorithms, low-precision software, systems, and hardware",
            "Changing the Cost: The empirical scaling law dictates what logical compute buys; the compounding efficiency stack constantly changes what that compute costs"
        ]
    },
    {
        "title": "Paper 4: Modernizing Amdahl’s Law",
        "bullets": [
            "Beyond Core Counts: Classical Amdahl’s Law (fixed serial/parallel fraction) fails to describe modern heterogeneous AI architectures",
            "Constrained Hardware Allocation: The modern design question is how to allocate a constrained hardware budget between specialized logic and programmable compute",
            "The Value-Scalable Fraction (S): The portion of the workload where injecting additional compute continues to create meaningful, end-to-end task value"
        ]
    },
    {
        "title": "Paper 4: The Collapse Threshold of Specialization",
        "bullets": [
            "[Visual Concept: A phase-boundary graph showing the race between Specialization Efficiency (R) and Scalable Fraction (S)]",
            "The Rising Bar: As the value-scalable workload (S) grows, fixed-function specialization faces a strict collapse threshold",
            "The Key Equation: For a given scalable fraction S, specialized hardware must achieve a relative efficiency ratio of R_c = 1/(1−S) to justify its inclusion",
            "Specialization must constantly re-earn its silicon"
        ]
    },
    {
        "title": "Architectural Consequence: The Pressure for Programmability",
        "bullets": [
            "Shifting Value: AI scaling pushes value-generation away from early fixed pipelines toward learned, late-stage computation (e.g., neural reconstruction, massive rerankers)",
            "Why GPUs Become More Programmable: As S rises, the optimal hardware locus shifts toward generalized, tensor-heavy programmable compute",
            "The Accelerator Dynamic: AI domain-specific accelerators matter, but they do not simply displace the GPU; they must continually adapt as the dominant value-producing workload evolves"
        ]
    },
    {
        "title": "Final Synthesis: The Unified Picture",
        "bullets": [
            "[Visual Concept: A single clean structural diagram linking formal limits, bounded feasibility, empirical scaling, and hardware architecture]",
            "1. Computability: Open-ended ascent across relative boundaries",
            "2. Feasibility: Staged reachability through polynomial-time verification",
            "3. Empirical Drivers: Logical scaling sustained by compounding hidden efficiency",
            "4. Architectural Response: Rising scalable workloads driving hardware toward programmable compute"
        ]
    },
    {
        "title": "Closing Takeaways",
        "bullets": [
            "A Structured Phenomenon: AI progress is neither pure hype nor a contradiction of theoretical limits",
            "The Unified Program: Intelligence emerges as a structured, open-ended process bounded by formal limits, propelled by empirical scaling drivers, and fundamentally reshaping system-level hardware consequences"
        ]
    }
]

# Theme Colors
BG_COLOR = RGBColor(10, 25, 41)       # Dark Navy
TEXT_COLOR = RGBColor(240, 240, 240)  # Off-white
ACCENT_CYAN = RGBColor(0, 255, 255)   # Blueprint Cyan
ACCENT_GOLD = RGBColor(255, 215, 0)   # Gold/Orange

def apply_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_text_box(slide, text, left, top, width, height, font_size=18, color=TEXT_COLOR, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial" 
    return txBox

def add_styled_shape(slide, shape_type, left, top, width, height, text="", font_size=16, line_col=ACCENT_CYAN, fill_col=BG_COLOR, text_col=TEXT_COLOR):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_col
    shape.line.color.rgb = line_col
    shape.line.width = Pt(2)
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_col
        p.font.bold = True
        p.font.name = "Arial"
    return shape

def draw_visual_concept(slide, concept_text, current_y):
    """Draws native PowerPoint shapes based on the specific visual concept."""
    y = current_y
    concept = concept_text.lower()
    
    if "blueprint-style diagram" in concept:
        # Mechanics of Ascent
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4), Inches(y+2), Inches(5), Inches(0.8), "A_n (Foundational Layer)")
        add_styled_shape(slide, MSO_SHAPE.UP_ARROW, Inches(6), Inches(y+1.0), Inches(1), Inches(0.8), "", line_col=ACCENT_GOLD, fill_col=ACCENT_GOLD)
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4), Inches(y), Inches(5), Inches(0.8), "Intermediate Step (B)")
        add_styled_shape(slide, MSO_SHAPE.UP_ARROW, Inches(6), Inches(y-0.8), Inches(1), Inches(0.6), "", line_col=ACCENT_GOLD, fill_col=ACCENT_GOLD)
        
        # Open Ceiling A'
        ceiling = add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(y-1.5), Inches(6), Inches(0.5), "A'_n (Turing Jump / Open Ceiling)", line_col=ACCENT_CYAN, text_col=ACCENT_CYAN)
        ceiling.line.dash_style = 4 # Dashed line
        return y + 3.0
        
    elif "finite-depth ladder" in concept:
        # Polynomial Hierarchy Ladder
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(y+1.6), Inches(3), Inches(0.6), "P (Direct Decision)", line_col=ACCENT_CYAN)
        add_styled_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(4.7), Inches(y+1.7), Inches(0.6), Inches(0.4), line_col=ACCENT_GOLD, fill_col=ACCENT_GOLD)
        
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(y+0.8), Inches(3.5), Inches(0.6), "NP (Certificate Verification)", line_col=ACCENT_CYAN)
        add_styled_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(9.2), Inches(y+0.9), Inches(0.6), Inches(0.4), line_col=ACCENT_GOLD, fill_col=ACCENT_GOLD)
        
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(10), Inches(y), Inches(3), Inches(0.6), "PH (Guided Search)", line_col=ACCENT_CYAN)
        return y + 2.5
        
    elif "nested boundary map" in concept:
        # Nested Bounds (Approximated with stacked rounded rectangles)
        add_styled_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(y), Inches(9.33), Inches(2.5), "Well-defined tasks (Open-ended)", line_col=ACCENT_CYAN)
        add_text_box(slide, "Well-defined tasks (Open-ended)", Inches(2.2), Inches(y+0.1), Inches(4), Inches(0.5), 18, ACCENT_CYAN, True)
        
        # We need to wipe text from the shape itself since we added a text box
        # Actually add_styled_shape adds text to the center. Let's recreate without text.
        pass # Re-drawing cleanly:
        s1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(y), Inches(9.33), Inches(2.5))
        s1.fill.solid(); s1.fill.fore_color.rgb = BG_COLOR; s1.line.color.rgb = ACCENT_CYAN
        add_text_box(slide, "Well-Defined Tasks (Open-Ended)", Inches(2.2), Inches(y+0.1), Inches(8), Inches(0.5), 16, ACCENT_CYAN, True)
        
        s2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(y+0.5), Inches(7.33), Inches(1.8))
        s2.fill.solid(); s2.fill.fore_color.rgb = BG_COLOR; s2.line.color.rgb = ACCENT_GOLD
        add_text_box(slide, "Computable Tasks", Inches(3.2), Inches(y+0.6), Inches(6), Inches(0.5), 16, ACCENT_GOLD, True)
        
        s3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(y+1.0), Inches(5.33), Inches(1.1))
        s3.fill.solid(); s3.fill.fore_color.rgb = BG_COLOR; s3.line.color.rgb = TEXT_COLOR
        add_text_box(slide, "Polynomial-Time Feasible Tasks", Inches(4.2), Inches(y+1.4), Inches(5), Inches(0.5), 16, TEXT_COLOR, True, PP_ALIGN.CENTER)
        
        return y + 2.8
        
    elif "pressure-conversion loop" in concept:
        # Cycle Diagram
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(2), Inches(y+0.5), Inches(3), Inches(1.0), "Logical Compute Demand", line_col=ACCENT_CYAN)
        add_styled_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(5.2), Inches(y+0.8), Inches(0.8), Inches(0.4), line_col=ACCENT_GOLD, fill_col=ACCENT_GOLD)
        
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(y+0.5), Inches(2.5), Inches(1.0), "Physical Burden", line_col=ACCENT_CYAN)
        add_styled_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(8.9), Inches(y+0.8), Inches(0.8), Inches(0.4), line_col=ACCENT_GOLD, fill_col=ACCENT_GOLD)
        
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(9.9), Inches(y+0.5), Inches(2.5), Inches(1.0), "Efficiency Optimization", line_col=ACCENT_CYAN)
        
        # Return loop underneath
        l1 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.5), Inches(y+1.8), Inches(7.5), Inches(0.3))
        l1.fill.solid(); l1.fill.fore_color.rgb = ACCENT_CYAN; l1.line.color.rgb = ACCENT_CYAN
        
        return y + 2.5
        
    elif "phase-boundary graph" in concept:
        # 2D Graph
        # Y Axis
        slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(2), Inches(y), Inches(0.1), Inches(2.5)).fill.solid()
        add_text_box(slide, "Specialization Efficiency (R)", Inches(0.1), Inches(y+0.5), Inches(2), Inches(1.0), 14, ACCENT_GOLD)
        # X Axis
        slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2), Inches(y+2.5), Inches(6), Inches(0.1)).fill.solid()
        add_text_box(slide, "Scalable Fraction (S)", Inches(4), Inches(y+2.6), Inches(3), Inches(0.5), 14, ACCENT_GOLD)
        
        # Curve
        # We can simulate a curve with a diagonal line
        line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(2), Inches(y+2.5), Inches(5), Inches(y+0.2))
        line.line.color.rgb = ACCENT_CYAN
        line.line.width = Pt(4)
        
        add_text_box(slide, "Collapse Threshold: R_c = 1 / (1-S)", Inches(5), Inches(y+0.5), Inches(4), Inches(1.0), 18, ACCENT_CYAN, True)
        return y + 3.0
        
    elif "single clean structural diagram" in concept:
        # Four connected boxes
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(y), Inches(2.5), Inches(1), "1. Formal Limits", line_col=ACCENT_GOLD)
        add_styled_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(3.6), Inches(y+0.4), Inches(0.4), Inches(0.2), fill_col=TEXT_COLOR, line_col=TEXT_COLOR)
        
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4.1), Inches(y), Inches(2.5), Inches(1), "2. Bounded Feasibility", line_col=ACCENT_CYAN)
        add_styled_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(6.7), Inches(y+0.4), Inches(0.4), Inches(0.2), fill_col=TEXT_COLOR, line_col=TEXT_COLOR)
        
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(y), Inches(2.5), Inches(1), "3. Empirical Scaling", line_col=ACCENT_GOLD)
        add_styled_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(9.8), Inches(y+0.4), Inches(0.4), Inches(0.2), fill_col=TEXT_COLOR, line_col=TEXT_COLOR)
        
        add_styled_shape(slide, MSO_SHAPE.RECTANGLE, Inches(10.3), Inches(y), Inches(2.5), Inches(1), "4. Hardware Architecture", line_col=ACCENT_CYAN)
        return y + 1.5

    # Default fallback if not matched
    add_text_box(slide, concept_text, Inches(1.5), Inches(y), Inches(10), Inches(1.0), 20, ACCENT_GOLD, True)
    return y + 1.5


def create_ai_blueprint_pptx(output_filename):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. SPECIAL SLIDE: TITLE SLIDE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_bg(slide)
    
    add_text_box(slide, slides_data[0]["title"], Inches(1), Inches(2.0), Inches(11.33), Inches(1.5), 54, ACCENT_CYAN, True, PP_ALIGN.CENTER)
    add_text_box(slide, slides_data[0]["bullets"][0], Inches(1), Inches(4.0), Inches(11.33), Inches(1), 24, TEXT_COLOR, False, PP_ALIGN.CENTER)
    
    # 2. ITERATE OVER REMAINING SLIDES
    for data in slides_data[1:]:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        apply_dark_bg(slide)
        
        add_text_box(slide, data["title"], Inches(0.8), Inches(0.5), Inches(11.73), Inches(1.0), 36, ACCENT_CYAN, True)
        
        current_y = 1.8 
        
        for bullet in data["bullets"]:
            bullet = bullet.strip()
            
            # Sub-routing for Visual Concepts
            if bullet.startswith("[Visual"):
                current_y = draw_visual_concept(slide, bullet, current_y)
                continue
                
            char_count = len(bullet)
            estimated_lines = max(1, (char_count // 70) + 1)
            height_needed = estimated_lines * 0.45 + 0.1
            
            font_size = 24
            color = TEXT_COLOR
            bold = False
            margin_left = 1.0
            
            if bullet.startswith("Pair") or "Connected Pairs" in bullet:
                color = ACCENT_GOLD
                font_size = 26
                bold = True
                margin_left = 0.8
                height_needed += 0.2
            elif bullet.startswith("Paper"):
                margin_left = 1.5
                font_size = 22
            else:
                margin_left = 1.2
            
            if ":" in bullet:
                height_needed += 0.1
                
            stripped_bullet = bullet.replace("Pair 1:", "").replace("Pair 2:", "") if "Pair" in bullet else bullet
            text_str = "• " + stripped_bullet if not stripped_bullet.startswith("  ") else stripped_bullet
                
            add_text_box(
                slide=slide, 
                text=text_str,
                left=Inches(margin_left), 
                top=Inches(current_y), 
                width=Inches(12.33 - margin_left), 
                height=Inches(height_needed), 
                font_size=font_size, 
                color=color, 
                bold=bold
            )
            
            current_y += height_needed + 0.1

    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}")

if __name__ == "__main__":
    create_ai_blueprint_pptx("The_AI_Progress_Blueprint.pptx")
